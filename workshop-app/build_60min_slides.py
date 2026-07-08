#!/usr/bin/env python3
"""Build 60-minute workshop slide deck (PPTX → import into Google Slides)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "60-min-workshop-slides.pptx"

BG = RGBColor(15, 17, 23)
ACCENT = RGBColor(91, 141, 255)
YOU = RGBColor(56, 189, 248)
HERMES = RGBColor(167, 139, 250)
TEXT = RGBColor(230, 232, 240)
MUTED = RGBColor(160, 165, 180)
CODE = RGBColor(200, 210, 230)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 24,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def _bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str = "",
    footer: str = "",
    title_color: RGBColor = ACCENT,
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide)
    _add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9), title, size=34, bold=True, color=title_color)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5), subtitle, size=16, color=MUTED)
    top = Inches(1.85) if subtitle else Inches(1.35)
    box = slide.shapes.add_textbox(Inches(0.75), top, Inches(11.5), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
    if footer:
        _add_textbox(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4), footer, size=14, color=MUTED)


def _code_slide(
    prs: Presentation,
    title: str,
    code: str,
    *,
    subtitle: str = "",
    title_color: RGBColor = HERMES,
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide)
    _add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8), title, size=32, bold=True, color=title_color)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.45), subtitle, size=15, color=MUTED)
    top = Inches(1.65) if subtitle else Inches(1.25)
    box = slide.shapes.add_textbox(Inches(0.65), top, Inches(12), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Menlo"
    p.font.color.rgb = CODE


def _two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
    *,
    left_color: RGBColor = YOU,
    right_color: RGBColor = HERMES,
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide)
    _add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8), title, size=32, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.6), Inches(1.35), Inches(5.8), Inches(0.5), left_title, size=22, bold=True, color=left_color)
    box_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.8))
    tf = box_l.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
    _add_textbox(slide, Inches(6.6), Inches(1.35), Inches(5.8), Inches(0.5), right_title, size=22, bold=True, color=right_color)
    box_r = slide.shapes.add_textbox(Inches(6.6), Inches(1.9), Inches(5.8), Inches(4.8))
    tf = box_r.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT


# Copy-paste blocks (also on slides)
INSTALL_CMD = """cd workshop-app/hermes-desktop
chmod +x install.sh
./install.sh
# Skills → ~/.hermes/skills/ — restart Hermes"""

MCP_CONFIG = """You.com MCP Server:
https://api.you.com/mcp

Enable tools on /account-action-brief:
• you-search
• you-contents
• you-research

(Optional) you-finance for public companies"""

TEST_CMD = """/account-action-brief

company: AMD
website: https://www.amd.com
goal: Prepare for an upcoming customer meeting — understand recent strategic signals and decide what actions the account team should consider.
output_audience: internal"""

DISCORD_INTRO = """Hi! I joined from the Account Action Brief workshop (Hermes + You.com).

- Name / role:
- What I'm building: /account-action-brief workflow
- First company I'll brief:
- One thing I want to learn in this community:"""


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    slide = _blank_slide(prs)
    _fill_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2), "Account Action Brief Workshop", size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.8), "60-Minute Live Track", size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.6), "You.com = evidence  ·  Hermes = governed workflow", size=20, color=MUTED, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5), "Slides = script + copy-paste  ·  App = govern + agent lab", size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # 2 Story
    _two_column_slide(
        prs,
        "Two platforms, one workflow",
        "You.com — evidence",
        [
            "Live web intelligence (not stale training data)",
            "Search → Contents → Research chain",
            "Citations + freshness in API responses",
            "MCP exposes tools to Hermes",
        ],
        "Hermes — workflow",
        [
            "Named skill: /account-action-brief",
            "Stable inputs + six-section output",
            "Connector map + review gates (Draft first)",
            "You install + test live today",
        ],
    )

    # 3 Agenda
    _bullet_slide(
        prs,
        "Agenda (60 min)",
        [
            "0–3   Welcome — story + demo company AMD",
            "3–18  You.com APIs — walk platform playground together",
            "18–23 Govern — connectors + review gates (app)",
            "23–26 Workflow card — generate + download (app)",
            "26–52 Hermes desktop — install + define inputs + MCP + test brief",
            "52–57 Close — join You.com Discord",
            "57–60 Buffer / Q&A",
        ],
        footer="Playground: https://you.com/platform",
    )

    # 4 API decision matrix
    _bullet_slide(
        prs,
        "You.com APIs — which one when?",
        [
            "Search — discover candidate sources (structured JSON, snippets, URLs)",
            "Contents — read full page text from URLs you already trust",
            "Research — multi-source synthesis with citations (output.content + sources)",
            "Finance Research — earnings, filings, market signals (public companies)",
            "MCP — Hermes calls Search/Contents/Research without custom integration code",
        ],
        subtitle="One job per API — Hermes orchestrates the chain",
        title_color=YOU,
    )

    # 5 Search
    _bullet_slide(
        prs,
        "Web Search API",
        [
            "Job: discover current web + news sources",
            "Returns: titles, URLs, snippets, freshness metadata",
            "Governance: include_domains, freshness, count",
            "Playground: try query → AMD data center Instinct MI350 earnings",
            "Point out: structured JSON — not a prose answer",
        ],
        subtitle="Playground → Search endpoint",
        title_color=YOU,
    )

    # 6 Contents
    _bullet_slide(
        prs,
        "Contents API",
        [
            "Job: fetch clean markdown/HTML from selected URLs",
            "Use when: snippets are not enough for claim verification",
            "Returns: page markdown + metadata per URL",
            "Playground: paste 1–2 URLs from Search results",
            "Point out: page-level evidence for the brief",
        ],
        subtitle="Playground → Contents endpoint",
        title_color=YOU,
    )

    # 7 Research
    _bullet_slide(
        prs,
        "Research API",
        [
            "Job: multi-step synthesis across sources",
            "Returns: output.content + output.sources (citations)",
            "Use when: you need strategic implications, not just retrieval",
            "Playground: AMD demo — attendees define their own goal in Hermes later",
            "Say: Hermes automates this same chain on /account-action-brief",
        ],
        subtitle="Playground → Research endpoint",
        title_color=YOU,
    )

    # 8 What is MCP + You.com MCP
    _two_column_slide(
        prs,
        "MCP — how Hermes calls You.com",
        "What is MCP?",
        [
            "Model Context Protocol — open standard for agent tool access",
            "Agents call tools (APIs) without custom integration code per app",
            "Playground = you call APIs manually; MCP = Hermes calls them as tools",
            "Same evidence chain: Search → Contents → Research",
        ],
        "You.com MCP server",
        [
            "URL: https://api.you.com/mcp",
            "Tools: you-search · you-contents · you-research",
            "Optional: finance research for public companies",
            "Your API key authenticates — read-only web intelligence",
        ],
    )

    # 9 MCP flow (full-width bullets)
    _bullet_slide(
        prs,
        "How You.com MCP works in /account-action-brief",
        [
            "1. You run /account-action-brief with company + brief goal",
            "2. Hermes skill decides which tools to call and in what order",
            "3. Hermes → MCP → you-search (discover sources for the goal)",
            "4. Hermes → MCP → you-contents (read selected URLs)",
            "5. Hermes → MCP → you-research (synthesize with citations)",
            "6. Skill formats six-section brief → review status: Draft",
            "You saw steps 3–5 manually in the playground — MCP automates them",
        ],
        subtitle="We wire https://api.you.com/mcp live during Hermes setup",
        title_color=YOU,
    )

    # 10 Finance optional
    _bullet_slide(
        prs,
        "Finance Research API (optional)",
        [
            "Finance-optimized index — earnings, filings, market signals",
            "Great for AMD and other public companies",
            "Often slow (2–5 min) — demo in playground if time allows",
            "research_effort: deep or exhaustive only",
            "Enable via MCP only if you need financial context in the brief",
        ],
        subtitle="Skip in 60-min if running long",
        title_color=YOU,
    )

    # 11 Playground walkthrough
    _bullet_slide(
        prs,
        "Platform playground walkthrough",
        [
            "1. Open https://you.com/platform → API Playground",
            "2. Search — run for AMD; show JSON results",
            "3. Contents — fetch one press release URL",
            "4. Research — ask about AMD strategic signals (renewal / outbound / competitive)",
            "5. View Code — show curl/Python snippet (optional)",
            "6. Bridge: 'In Hermes, MCP runs this chain automatically'",
        ],
        subtitle="3–18 min · facilitator leads; attendees follow along",
        title_color=YOU,
    )

    # 12 Hermes skills + connectors
    _two_column_slide(
        prs,
        "Hermes — skills & connectors",
        "Skills (workflows)",
        [
            "Named, reusable workflow — e.g. /account-action-brief",
            "SKILL.md = inputs, procedure, output schema, governance rules",
            "Trigger: slash command + company + brief goal",
            "Today: install skill pack → open SKILL.md to see what's inside",
            "Skills define WHAT the agent does every run",
        ],
        "Connectors (permissions)",
        [
            "Read-only: You.com MCP, CRM read, internal docs",
            "Draft-only: email draft, Slack review, CRM update draft",
            "Blocked: auto-send, auto-CRM-update, sensitive data",
            "Connectors define HOW the skill may touch each system",
            "You.com is always read-only — Hermes never writes via MCP",
        ],
        left_color=HERMES,
        right_color=ACCENT,
    )

    # 13 Govern
    _two_column_slide(
        prs,
        "Govern — review gates",
        "Why governance?",
        [
            "AI output starts Draft — human review before use",
            "Internal use ≠ external use (customer comms need extra approval)",
            "Block auto-send / auto-CRM-update in pilot — humans approve",
            "Skills + connectors + review gates = safe, repeatable workflow",
        ],
        "Review gates",
        [
            "Draft — default every run",
            "Needs edits — weak sources or unsupported claims",
            "Approved for internal use — meeting prep, planning",
            "Approved for external use — outreach, customer comms",
        ],
    )

    # 14 Define inputs in Hermes
    _bullet_slide(
        prs,
        "Define inputs in Hermes (live)",
        [
            "Required: company — AMD (demo) or your account",
            "Optional: website — https://www.amd.com",
            "Brief goal — tailors Search + Research queries:",
            "  Renewal · Outbound · Competitive · Partner",
            "Optional: output_audience — internal or external draft",
            "Six-section output every brief follows:",
            "  Snapshot · Signals · Why They Care · Actions · Claims & Sources · Review Notes",
        ],
        subtitle="26–41 min · define in test command, then run /account-action-brief",
        title_color=HERMES,
    )

    # 15 Hermes install
    _code_slide(
        prs,
        "Hermes desktop — install skill pack (live)",
        INSTALL_CMD,
        subtitle="26–34 min · everyone runs this in the room",
    )

    # 16 MCP wire (copy-paste)
    _code_slide(
        prs,
        "Wire You.com MCP (live)",
        MCP_CONFIG,
        subtitle="34–41 min · circulate and help with tool permissions",
    )

    # 17 Test
    _code_slide(
        prs,
        "Define inputs + test /account-action-brief (live)",
        TEST_CMD,
        subtitle="41–52 min · company + goal + verify six sections + Draft",
    )

    # 18 Verification
    _bullet_slide(
        prs,
        "Minimum success before you leave",
        [
            "Skill /account-action-brief installed in Hermes",
            "You.com MCP connected (tool calls visible)",
            "Test run completes with six sections",
            "Claims & Sources lists real URLs",
            "Review status = Draft",
            "No auto-send or CRM write occurred",
        ],
        title_color=HERMES,
    )

    # 19 Discord
    _code_slide(
        prs,
        "Close — join You.com Discord",
        f"https://discord.gg/2C4WgryxSD\n\nPost in #introductions:\n\n{DISCORD_INTRO}",
        subtitle="52–57 min · only required follow-up",
        title_color=ACCENT,
    )

    # 20 Close script
    slide = _blank_slide(prs)
    _fill_bg(slide)
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8), "Close script (read aloud)", size=36, bold=True, color=ACCENT)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.2),
        Inches(3.8),
        (
            '"You walked the You.com APIs in the playground, installed Hermes, '
            "and ran a governed account brief.\n\n"
            'Join the You.com Discord — link on screen — introduce yourself '
            "and say you came from this workshop.\n\n"
            'Keep using /account-action-brief on your accounts. See you in the server."'
        ),
        size=22,
        color=TEXT,
    )

    # 21 Resources
    _bullet_slide(
        prs,
        "Resources",
        [
            "You.com Playground: https://you.com/platform",
            "Workshop app: http://localhost:8080",
            "FACILITATOR-60MIN.md · HERMES-DESKTOP-SETUP.md",
            "COMMUNITY-CONTINUATION.md (Discord only)",
            "Preflight: ./preflight.sh",
        ],
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
